"""The script stage as a dispatcher, and the per-topic overlays.

Every extractor emits the same per-topic structure, which is the whole reason
the aligner, the checks, the artifacts and the packet did not change when a
second source was added. These tests pin that shape, and the two states that
are applied on top of whichever extractor ran.

The deck is generated. Nothing here quotes a real storyboard.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from qa.extract_script import build_script_for_source
from qa.script_source import DOCX_BUS, FREEFORM, NONE, OUTLINE, PPTX, TopicScript
from qa.util import ScriptError

NARRATION = {
    "01": "The kettle boils a fresh pot. Volunteers gather in the greenhouse.",
    "02": "In this video, we repot the tallest cuttings. The soil is warm.",
    "03": "In this demonstration, we water the seedlings. Watch the gauge.",
}


def make_deck(path: Path, notes: dict[str, str] | None = None) -> Path:
    """A storyboard with one narrated slide per topic, plus a template slide."""
    from pptx import Presentation

    notes = notes if notes is not None else NARRATION
    deck = Presentation()
    boiler = deck.slides.add_slide(deck.slide_layouts[5])
    boiler.shapes.title.text = "Directions"
    boiler.notes_slide.notes_text_frame.text = (
        "General directions for using this template. Delete this slide."
    )
    for title, text in notes.items():
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = f"Topic {title}"
        slide.notes_slide.notes_text_frame.text = text
    deck.save(str(path))
    return path


@pytest.fixture
def deck(tmp_path) -> Path:
    return make_deck(tmp_path / "storyboard.pptx")


TOPICS = ["01", "02", "03"]

# The keys every extractor's per-topic entry has to carry, whatever it read.
REQUIRED = {
    "topic",
    "script",
    "scripted",
    "slides",
    "source_ref",
    "sentences",
    "non_narration",
    "word_count",
}


# ---------------------------------------------------------------------------
# One shape, whichever extractor ran
# ---------------------------------------------------------------------------

def test_the_pptx_extractor_emits_the_common_shape(deck, tmp_path):
    script = build_script_for_source(PPTX, deck, TOPICS, {}, tmp_path)
    assert script["script_source"] == PPTX
    assert script["script_document"] == "storyboard.pptx"
    for entry in script["topics"]:
        assert REQUIRED <= set(entry), entry["topic"]
        assert entry["script"] == "verbatim"
        assert entry["scripted"] is True
        assert entry["non_narration"] == []
        assert entry["source_ref"].startswith("slides ")


def test_the_docx_extractor_emits_the_same_shape(tmp_path):
    from test_bus_template import make_bus_document

    document = make_bus_document(tmp_path / "scripts.docx")
    script = build_script_for_source(
        DOCX_BUS, document, TOPICS, {}, tmp_path, course_code="it_gen01_02_enus"
    )
    assert script["script_source"] == DOCX_BUS
    for entry in script["topics"]:
        assert REQUIRED <= set(entry), entry["topic"]
        assert entry["slides"] is None, "a Word script has no slides"
        assert entry["source_ref"]


def test_a_source_with_no_extractor_is_refused(deck, tmp_path):
    with pytest.raises(ScriptError, match="No extractor"):
        build_script_for_source(FREEFORM, deck, TOPICS, {}, tmp_path)


# ---------------------------------------------------------------------------
# The overlays
# ---------------------------------------------------------------------------

def test_an_outline_topic_keeps_its_text_and_is_not_aligned(deck, tmp_path):
    states = {"03": TopicScript(state=OUTLINE)}
    script = build_script_for_source(PPTX, deck, TOPICS, states, tmp_path)
    demo = next(t for t in script["topics"] if t["topic"] == "03")
    assert demo["script"] == OUTLINE
    assert demo["scripted"] is False
    assert demo["sentences"], "the outline text is still carried"
    assert demo["outline"], "and kept separately, for the packet to show"


def test_a_none_topic_keeps_its_place_and_carries_no_script(deck, tmp_path):
    """Not skipped. The audio was delivered and still has to be checked."""
    states = {"03": TopicScript(state=NONE)}
    script = build_script_for_source(PPTX, deck, TOPICS, states, tmp_path)

    assert [t["topic"] for t in script["topics"]] == TOPICS
    demo = next(t for t in script["topics"] if t["topic"] == "03")
    assert demo["script"] == NONE
    assert demo["scripted"] is False
    assert demo["sentences"] == []
    assert demo["word_count"] == 0
    assert "outline" not in demo
    # It still says where in the deck it sits, which is a real fact about it.
    assert demo["slides"]


def test_a_freeform_topic_reads_its_own_document(deck, tmp_path):
    (tmp_path / "demo_script.txt").write_text(
        "Open the console. Choose the second option.", encoding="utf-8"
    )
    states = {"03": TopicScript(state=FREEFORM, file="demo_script.txt")}
    script = build_script_for_source(PPTX, deck, TOPICS, states, tmp_path)

    demo = next(t for t in script["topics"] if t["topic"] == "03")
    assert demo["script"] == FREEFORM
    assert demo["scripted"] is True
    assert demo["sentences"] == [
        "Open the console.",
        "Choose the second option.",
    ]
    assert demo["source_ref"] == "demo_script.txt"


def test_a_freeform_topic_stops_pointing_at_the_slides_it_replaced(deck, tmp_path):
    """Its script is elsewhere now, so the slide reference would misdirect."""
    (tmp_path / "demo_script.txt").write_text("Open the console.", encoding="utf-8")
    states = {"03": TopicScript(state=FREEFORM, file="demo_script.txt")}
    script = build_script_for_source(PPTX, deck, TOPICS, states, tmp_path)

    demo = next(t for t in script["topics"] if t["topic"] == "03")
    assert demo["slides"] is None
    assert "In this demonstration" not in " ".join(demo["sentences"])


def test_a_missing_freeform_document_stops_the_run(deck, tmp_path):
    states = {"03": TopicScript(state=FREEFORM, file="not_there.txt")}
    with pytest.raises(ScriptError):
        build_script_for_source(PPTX, deck, TOPICS, states, tmp_path)


def test_the_overlays_reach_a_word_script_too(tmp_path):
    """Applied after whichever extractor ran, so both get them."""
    from test_bus_template import make_bus_document

    document = make_bus_document(tmp_path / "scripts.docx")
    states = {"03": TopicScript(state=NONE)}
    script = build_script_for_source(
        DOCX_BUS, document, TOPICS, states, tmp_path, course_code="it_gen01_02_enus"
    )
    summary = next(t for t in script["topics"] if t["topic"] == "03")
    assert summary["script"] == NONE
    assert summary["sentences"] == []


# ---------------------------------------------------------------------------
# The legacy key still means what it always meant
# ---------------------------------------------------------------------------

def test_unscripted_topics_and_topic_states_say_the_same_thing(deck, tmp_path):
    from qa.extract_script import build_script

    by_legacy = build_script(deck, TOPICS, unscripted={"03"})
    by_state = build_script(
        deck, TOPICS, topic_scripts={"03": TopicScript(state=OUTLINE)}
    )
    assert by_legacy["topics"] == by_state["topics"]


def test_an_explicit_state_wins_over_the_legacy_key(deck, tmp_path):
    """Someone who wrote out a state meant it."""
    from qa.extract_script import build_script

    script = build_script(
        deck,
        TOPICS,
        unscripted={"03"},
        topic_scripts={"03": TopicScript(state=NONE)},
    )
    demo = next(t for t in script["topics"] if t["topic"] == "03")
    assert demo["script"] == NONE


def test_an_overlay_does_not_put_one_extractors_key_into_the_others_output(
    deck, tmp_path
):
    """A storyboard's sentences come from slides; a Word script's from rows.

    Clearing one must not create the other, or an entry ends up claiming a
    provenance its extractor never had.
    """
    from test_bus_template import make_bus_document

    states = {"03": TopicScript(state=NONE)}

    from_deck = build_script_for_source(PPTX, deck, TOPICS, states, tmp_path)
    demo = next(t for t in from_deck["topics"] if t["topic"] == "03")
    assert demo["sentence_slides"] == []
    assert "sentence_rows" not in demo

    document = make_bus_document(tmp_path / "scripts.docx")
    from_word = build_script_for_source(
        DOCX_BUS, document, TOPICS, states, tmp_path, course_code="it_gen01_02_enus"
    )
    summary = next(t for t in from_word["topics"] if t["topic"] == "03")
    assert summary["sentence_rows"] == []
    assert "sentence_slides" not in summary


# ---------------------------------------------------------------------------
# The source is defaulted from the project type, not welded to it
# ---------------------------------------------------------------------------

def plant_course(course: Path, project_type: str, document: str, source: str = "") -> Path:
    course.mkdir(parents=True, exist_ok=True)
    (course / document).write_bytes(b"")
    lines = [
        'course_number: "02"',
        f"project_type: {project_type}",
        "course_code: it_gen01_02_enus",
    ]
    if source:
        lines.append(f"script_source: {source}")
    (course / "course.yaml").write_text("\n".join(lines) + "\n", encoding="utf-8")
    return course


def test_the_source_defaults_from_the_project_type(tmp_path):
    from qa.config import load_course_yaml

    vendor = plant_course(tmp_path / "v", "VENDOR", "deck.pptx")
    assert load_course_yaml(vendor).script_source == PPTX

    cgt = plant_course(tmp_path / "c", "CGT", "scripts.docx")
    assert load_course_yaml(cgt).script_source == DOCX_BUS


def test_an_unusual_course_may_state_a_source_its_type_does_not_imply(tmp_path):
    """The default is a default. A CGT course with a deck is allowed to say so."""
    from qa.config import load_course_yaml

    course = plant_course(tmp_path / "odd", "CGT", "deck.pptx", source="pptx")
    cfg = load_course_yaml(course)
    assert cfg.script_source == PPTX
    assert cfg.script_document.name == "deck.pptx"


def test_a_per_topic_state_may_not_be_claimed_as_the_courses_source(tmp_path):
    """A course whose every topic is unscripted is not a course this can read."""
    from qa.config import load_course_yaml
    from qa.util import ConfigError

    course = plant_course(tmp_path / "bad", "VENDOR", "deck.pptx", source="none")
    with pytest.raises(ConfigError, match="describes one topic, not a"):
        load_course_yaml(course)


def test_an_unknown_source_is_refused_by_name(tmp_path):
    from qa.config import load_course_yaml
    from qa.util import ConfigError

    course = plant_course(tmp_path / "huh", "VENDOR", "deck.pptx", source="markdown")
    with pytest.raises(ConfigError, match="is not recognized"):
        load_course_yaml(course)
