"""Tests for the environment check and setup command.

Every prerequisite is exercised against injected results rather than against
the machine the tests happen to run on, because a check that only passes on a
correctly configured desktop is not a check.

The rule the command exists to keep, from D2: it never installs system
software. Several tests here assert that directly.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from qa.device import Device
from qa.setup import (
    CTRANSLATE2_CUDA_MAJOR,
    MAX_PYTHON_EXCLUSIVE,
    MISMATCH,
    MISSING,
    NOT_USABLE,
    OK,
    Requirement,
    SetupReport,
    build_fixture,
    check_cuda,
    check_model,
    check_packages,
    check_python,
    check_tool,
    check_venv,
    cuda_pip_remediation,
    download_model,
    install_hint,
    install_packages,
    main,
    render_fixes,
    render_table,
    smoke_test,
)

USABLE_GPU = Device("cuda", "GPU", True, detail="1 CUDA device, compute types int8")
NO_GPU = Device(
    "cuda", "GPU", False, reason="no CUDA capable device was found. This machine has no NVIDIA GPU"
)
BROKEN_GPU = Device(
    "cuda", "GPU", False, reason="a CUDA device was found but the runtime cannot use it"
)


# ---------------------------------------------------------------------------
# Python, and the newer-is-worse trap
# ---------------------------------------------------------------------------

def test_a_supported_python_is_ok():
    assert check_python((3, 12, 10)).status == OK
    assert check_python((3, 11, 0)).status == OK


def test_too_old_a_python_is_a_mismatch_with_a_command():
    found = check_python((3, 9, 18))
    assert found.status == MISMATCH
    assert found.fix, "a mismatch must come with the command that fixes it"


def test_a_newer_python_is_a_trap_not_an_upgrade():
    """D9: the ASR runtime publishes no wheels for it, so it fails to build."""
    found = check_python((MAX_PYTHON_EXCLUSIVE[0], MAX_PYTHON_EXCLUSIVE[1], 0))
    assert found.status == MISMATCH
    assert "D9" in found.detail
    assert "trap" in found.detail.lower()
    assert found.fix


def test_python_is_never_reported_as_merely_missing():
    """Something is running this command, so "not installed" is never the answer."""
    for version in ((3, 9, 0), (3, 14, 0), (3, 12, 0)):
        assert check_python(version).status != MISSING


# ---------------------------------------------------------------------------
# System binaries: found, missing, and always with a fix
# ---------------------------------------------------------------------------

def test_a_missing_tool_reports_the_platform_command(monkeypatch):
    monkeypatch.setattr("qa.setup.shutil.which", lambda name: None)
    found = check_tool("ffmpeg")
    assert found.status == MISSING
    assert found.fix == install_hint("ffmpeg")
    assert found.fix, "a missing tool with no remediation is a dead end"


def test_a_present_tool_reports_its_version(monkeypatch):
    monkeypatch.setattr("qa.setup.shutil.which", lambda name: f"/usr/bin/{name}")
    monkeypatch.setattr("qa.setup._version_of", lambda command: "ffmpeg version 9.0.1")
    found = check_tool("ffmpeg")
    assert found.status == OK
    assert "9.0.1" in found.found


@pytest.mark.parametrize("system,needle", [
    ("Windows", "winget"), ("Linux", "apt-get"), ("Darwin", "brew"),
])
def test_remediation_is_written_for_the_platform_in_hand(monkeypatch, system, needle):
    monkeypatch.setattr("qa.setup._system", lambda: system)
    assert needle in install_hint("ffmpeg")
    assert needle in install_hint("git")


def test_an_unknown_platform_still_gets_somewhere_to_go(monkeypatch):
    monkeypatch.setattr("qa.setup._system", lambda: "Haiku")
    assert "ffmpeg.org" in install_hint("ffmpeg")


# ---------------------------------------------------------------------------
# CUDA: three states that must not collapse into one
# ---------------------------------------------------------------------------

def test_a_usable_gpu_is_ok():
    found = check_cuda(probe=USABLE_GPU, toolkit=(None, ""))
    assert found.status == OK
    assert found.optional, "GPU must never block setup"


def test_no_gpu_at_all_is_missing_not_a_mismatch():
    found = check_cuda(probe=NO_GPU, toolkit=(None, ""))
    assert found.status == MISSING
    assert found.optional
    assert "runs on CPU" in found.fix


def test_an_old_toolkit_with_an_unusable_gpu_is_a_version_mismatch():
    """The live case this row exists for: a real card the engine cannot use."""
    found = check_cuda(probe=BROKEN_GPU, toolkit=(11, "on PATH: C:/CUDA/v11.0/bin"))
    assert found.status == MISMATCH
    assert found.status != MISSING, "an old toolkit is not an absence"
    assert "CUDA 11" in found.found
    assert str(CTRANSLATE2_CUDA_MAJOR) in found.required


def test_the_mismatch_remediation_is_additive_and_local():
    found = check_cuda(probe=BROKEN_GPU, toolkit=(11, "on PATH: C:/CUDA/v11.0/bin"))
    assert "pip install" in found.fix
    assert "cu12" in found.fix
    assert "virtual environment" in found.detail
    assert "nothing system wide" in found.detail


def test_the_mismatch_remediation_never_says_to_uninstall():
    """CUDA versions coexist, and something else may depend on the old one."""
    found = check_cuda(probe=BROKEN_GPU, toolkit=(11, "on PATH: C:/CUDA/v11.0"))
    text = (found.fix + " " + found.detail).lower()

    # Every mention of uninstalling must be a prohibition, not advice. A bare
    # substring ban would fail on the sentence that does the prohibiting, so
    # the property to assert is that the two counts match.
    assert text.count("uninstall") > 0, "it should address the temptation"
    assert text.count("uninstall") == text.count("do not uninstall")
    assert "coexist" in text, "the reason matters as much as the instruction"

    # Removing an orphaned toolkit may be mentioned as the user's own call,
    # never as a step of this remediation.
    if "remov" in text:
        assert "your call" in text


def test_a_broken_gpu_with_no_old_toolkit_is_not_usable():
    found = check_cuda(probe=BROKEN_GPU, toolkit=(None, ""))
    assert found.status == NOT_USABLE
    assert found.optional


def test_a_usable_gpu_alongside_an_old_toolkit_still_reads_ok():
    """What matters is whether the engine can use it, not what else is on disk."""
    found = check_cuda(probe=USABLE_GPU, toolkit=(11, "on PATH: C:/CUDA/v11.0"))
    assert found.status == OK
    assert "CUDA 11" in found.found, "the old toolkit is still worth mentioning"


def test_gpu_is_optional_in_every_state():
    for probe in (USABLE_GPU, NO_GPU, BROKEN_GPU):
        assert check_cuda(probe=probe, toolkit=(None, "")).optional


def test_the_pip_remediation_names_both_libraries():
    text = cuda_pip_remediation()
    assert "nvidia-cublas-cu12" in text and "nvidia-cudnn-cu12" in text


# ---------------------------------------------------------------------------
# The environment this project installs into
# ---------------------------------------------------------------------------

def test_a_missing_venv_is_reported_not_assumed(tmp_path):
    found = check_venv(tmp_path / "nothing")
    assert found.status == MISSING
    assert "will create it" in found.fix


def test_packages_cannot_be_checked_without_a_venv(tmp_path):
    found = check_packages(tmp_path / "nothing")
    assert found.status == MISSING


def test_installing_packages_is_skipped_when_already_satisfied(monkeypatch, tmp_path):
    """Idempotence: rerunning setup must not reinstall."""
    executable = tmp_path / "Scripts" / "python.exe"
    executable.parent.mkdir(parents=True)
    executable.write_text("", encoding="utf-8")
    monkeypatch.setattr("qa.setup.venv_python", lambda venv=None: executable)
    monkeypatch.setattr(
        "qa.setup.check_packages", lambda venv=None: Requirement("Python packages", OK)
    )

    def explode(*args, **kwargs):
        raise AssertionError("pip must not run when the packages are already there")

    monkeypatch.setattr("qa.setup.subprocess.run", explode)
    changed, message = install_packages(tmp_path)
    assert changed is False
    assert "already installed" in message


def test_creating_a_venv_is_skipped_when_one_exists(monkeypatch, tmp_path):
    from qa.setup import create_venv

    executable = tmp_path / "Scripts" / "python.exe"
    executable.parent.mkdir(parents=True)
    executable.write_text("", encoding="utf-8")
    monkeypatch.setattr("qa.setup.venv_python", lambda venv=None: executable)

    def explode(*args, **kwargs):
        raise AssertionError("venv must not be recreated")

    monkeypatch.setattr("qa.setup.subprocess.run", explode)
    changed, message = create_venv(tmp_path)
    assert changed is False
    assert "already present" in message


# ---------------------------------------------------------------------------
# The model: never a silent three gigabytes
# ---------------------------------------------------------------------------

def test_the_model_is_never_downloaded_without_confirmation(monkeypatch):
    monkeypatch.setattr(
        "qa.setup.check_model",
        lambda name="large-v3": Requirement(f"ASR model ({name})", MISSING),
    )

    def explode(*args, **kwargs):
        raise AssertionError("no download may start without a yes")

    monkeypatch.setattr("qa.setup.subprocess.run", explode)
    changed, message = download_model("large-v3", confirm=False)
    assert changed is False
    assert "not downloaded" in message
    assert "first transcription will download it" in message


def test_a_cached_model_is_not_downloaded_again(monkeypatch):
    monkeypatch.setattr(
        "qa.setup.check_model",
        lambda name="large-v3": Requirement(f"ASR model ({name})", OK),
    )

    def explode(*args, **kwargs):
        raise AssertionError("a cached model must not be fetched again")

    monkeypatch.setattr("qa.setup.subprocess.run", explode)
    changed, message = download_model("large-v3", confirm=True)
    assert changed is False
    assert "already cached" in message


def test_a_missing_model_does_not_block_setup():
    assert check_model("a-model-that-does-not-exist").optional


# ---------------------------------------------------------------------------
# --check changes nothing
# ---------------------------------------------------------------------------

def test_check_mode_performs_no_actions(monkeypatch, capsys):
    """The diagnostic doubles as the troubleshooting tool; it must be inert."""
    def explode(name):
        def inner(*args, **kwargs):
            raise AssertionError(f"--check must not call {name}")
        return inner

    monkeypatch.setattr("qa.setup.create_venv", explode("create_venv"))
    monkeypatch.setattr("qa.setup.install_packages", explode("install_packages"))
    monkeypatch.setattr("qa.setup.download_model", explode("download_model"))
    monkeypatch.setattr("qa.setup.smoke_test", explode("smoke_test"))

    main(["--check"])
    printed = capsys.readouterr().out
    assert "PREREQUISITE" in printed


def test_check_mode_writes_nothing_to_the_project(tmp_path, monkeypatch, capsys):
    before = {p: p.stat().st_mtime for p in Path("qa").glob("*.py")}
    main(["--check"])
    after = {p: p.stat().st_mtime for p in Path("qa").glob("*.py")}
    assert before == after


def test_check_mode_reports_failure_with_a_nonzero_exit(monkeypatch, capsys):
    monkeypatch.setattr(
        "qa.setup.check_all",
        lambda venv=None: [Requirement("ffmpeg", MISSING, "not on PATH", "any", "install it")],
    )
    assert main(["--check"]) == 1
    printed = capsys.readouterr().out
    assert "install it" in printed


def test_check_mode_reports_success_with_zero(monkeypatch, capsys):
    monkeypatch.setattr(
        "qa.setup.check_all", lambda venv=None: [Requirement("ffmpeg", OK, "9.0.1", "any")]
    )
    assert main(["--check"]) == 0


def test_setup_stops_rather_than_installing_system_software(monkeypatch, capsys):
    """D2, enforced: a missing system binary ends the run with instructions."""
    monkeypatch.setattr(
        "qa.setup.check_all",
        lambda venv=None: [
            Requirement("ffmpeg", MISSING, "not on PATH", "any", "winget install ..."),
        ],
    )

    def explode(*args, **kwargs):
        raise AssertionError("setup must not proceed past missing system software")

    monkeypatch.setattr("qa.setup.create_venv", explode)
    assert main([]) == 1
    printed = capsys.readouterr().out
    assert "does not install system software" in printed


# ---------------------------------------------------------------------------
# The table itself
# ---------------------------------------------------------------------------

def test_the_table_shows_all_four_states_distinctly():
    rows = [
        Requirement("a", OK, "1.0", "any"),
        Requirement("b", MISSING, "not on PATH", "any", "install b"),
        Requirement("c", MISMATCH, "v11", "v12", "upgrade c"),
        Requirement("d", NOT_USABLE, "present but broken", "usable", "fix d", optional=True),
    ]
    table = render_table(rows)
    for state in (OK, MISSING, MISMATCH, NOT_USABLE):
        assert state in table
    assert MISSING != MISMATCH


def test_the_fix_block_names_every_problem_and_its_command():
    rows = [
        Requirement("a", OK, "1.0", "any"),
        Requirement("b", MISSING, "gone", "any", "install b"),
        Requirement("c", MISMATCH, "v11", "v12", "upgrade c", optional=True),
    ]
    fixes = render_fixes(rows)
    assert "install b" in fixes and "upgrade c" in fixes
    assert "does not block setup" in fixes, "an optional failure must say so"
    assert "\n  a:" not in fixes, "a satisfied row needs no advice"


def test_a_clean_machine_prints_no_advice():
    assert render_fixes([Requirement("a", OK, "1.0", "any")]) == ""


def test_a_report_is_ok_when_only_optional_items_fail():
    report = SetupReport(
        requirements=[
            Requirement("ffmpeg", OK),
            Requirement("CUDA (GPU)", MISMATCH, optional=True),
        ]
    )
    assert report.ok
    assert report.blocking == []


# ---------------------------------------------------------------------------
# Proof: the fixture and the smoke test
# ---------------------------------------------------------------------------

def test_the_fixture_is_a_real_course(tmp_path):
    course = build_fixture(tmp_path)
    assert (course / "course.yaml").exists()
    assert list(course.glob("*.pptx"))
    assert list((course / "audio").glob("*.wav"))

    from qa.config import load_course_yaml

    config = load_course_yaml(course)
    assert config.course_code == "it_smoke_01_enus"
    assert config.project_type == "VENDOR"


def test_the_fixture_carries_no_course_narration(tmp_path):
    """D16: nothing in this repository quotes customer material."""
    course = build_fixture(tmp_path)
    from pptx import Presentation

    deck = Presentation(str(next(course.glob("*.pptx"))))
    notes = " ".join(
        slide.notes_slide.notes_text_frame.text
        for slide in deck.slides
        if slide.has_notes_slide
    ).lower()
    for word in ("cloud", "security", "siem", "storyboard narration"):
        assert word not in notes


def test_the_fixture_is_small(tmp_path):
    course = build_fixture(tmp_path)
    total = sum(p.stat().st_size for p in course.rglob("*") if p.is_file())
    assert total < 2_000_000, "the smoke fixture must stay small"


@pytest.mark.skipif(
    not check_model("tiny").satisfied and not check_model("large-v3").satisfied,
    reason="no ASR model cached; run qa-setup to fetch one",
)
def test_the_smoke_test_passes_end_to_end(tmp_path):
    """Installed and works are different claims. This is the second one."""
    passed, detail = smoke_test(tmp_path)
    assert passed, detail
    assert "reconciliation_packet" in detail


def test_the_smoke_test_says_so_when_no_model_is_available(monkeypatch, tmp_path):
    monkeypatch.setattr(
        "qa.setup.check_model",
        lambda name="large-v3": Requirement(f"ASR model ({name})", MISSING),
    )
    passed, detail = smoke_test(tmp_path)
    assert passed is False
    assert "skipped" in detail
