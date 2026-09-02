"""What a job record claims, against what is actually true.

A status record is a claim a process makes about itself, and a process that
dies stops updating its claim. Read literally, a killed run says RUNNING for
ever and every reader waits for something that is not coming. That is the same
failure as a transcriber truncating a file without saying so, which is the one
this codebase exists to prevent, so the record gets checked rather than
believed.

The bug these were written for: a run started from the web UI completed and
wrote its packet while the progress view said "Queued" throughout, and kept
saying it afterwards. The engine was right and the instrument reporting on the
engine was lying.
"""

from __future__ import annotations

import json
import os
import subprocess
import sys
import time
from dataclasses import dataclass
from pathlib import Path

import pytest

from qa.jobs import (
    DONE,
    FAILED,
    PENDING,
    RUNNING,
    FileJobStore,
    JobError,
    JobStatus,
    process_alive,
    resolve,
    submit,
    wrote_a_packet,
)


@pytest.fixture
def store(tmp_path) -> FileJobStore:
    return FileJobStore(tmp_path / "jobs")


def make_course(tmp_path: Path, name: str = "course10") -> Path:
    course = tmp_path / "library" / "spisccc26" / name
    course.mkdir(parents=True, exist_ok=True)
    (course / "course.yaml").write_text(
        'course_number: "10"\nproject_type: VENDOR\n'
        "course_code: it_spisccc26_10_enus\nreviewed_by: Ryan\n",
        encoding="utf-8",
    )
    return course


def wrote_packet(course: Path, when: float | None = None) -> Path:
    out = course / "qa_out"
    out.mkdir(parents=True, exist_ok=True)
    marker = out / "packet_index.json"
    marker.write_text(json.dumps({"path": str(out / "p.md")}), encoding="utf-8")
    if when is not None:
        os.utime(marker, (when, when))
    return marker


@dataclass
class FakeProcess:
    pid: int = 999999


# ---------------------------------------------------------------------------
# Asking the operating system
# ---------------------------------------------------------------------------

def test_this_process_is_alive_and_pid_zero_is_not():
    assert process_alive(os.getpid()) is True
    assert process_alive(0) is False


def test_a_process_that_has_exited_is_not_alive():
    child = subprocess.Popen([sys.executable, "-c", "pass"])
    child.wait(timeout=30)
    assert process_alive(child.pid) is False


def test_the_liveness_probe_does_not_kill_what_it_asks_about():
    """os.kill on Windows terminates. This must never reach for it."""
    child = subprocess.Popen([sys.executable, "-c", "import time; time.sleep(30)"])
    try:
        assert process_alive(child.pid) is True
        assert process_alive(child.pid) is True
        assert child.poll() is None, "the probe killed the process it asked about"
    finally:
        child.kill()
        child.wait(timeout=10)


# ---------------------------------------------------------------------------
# Reconciling a record with reality
# ---------------------------------------------------------------------------

def test_a_running_job_whose_process_lives_is_left_alone(store, tmp_path):
    job = JobStatus(
        id="live",
        course_dir=str(make_course(tmp_path)),
        state=RUNNING,
        pid=os.getpid(),
        started_at=time.time(),
    )
    store.write(job)
    assert resolve(store.read("live"), store).state == RUNNING


def test_a_dead_job_that_wrote_a_packet_is_finished(store, tmp_path):
    """The evidence a run finished does not have to come from the run."""
    course = make_course(tmp_path)
    started = time.time()
    wrote_packet(course, started + 5)

    job = JobStatus(
        id="lost", course_dir=str(course), state=PENDING, pid=999999, started_at=started
    )
    store.write(job)

    healed = resolve(store.read("lost"), store)
    assert healed.state == DONE
    assert healed.finished_at
    assert any("lost its last write" in w or "before it could record" in w
               for w in healed.warnings)


def test_a_dead_job_with_no_packet_is_failed_with_what_it_printed(store, tmp_path):
    course = make_course(tmp_path)
    job = JobStatus(
        id="died",
        course_dir=str(course),
        state=RUNNING,
        pid=999999,
        started_at=time.time(),
    )
    store.write(job)
    store.log_path("died").write_text(
        "Traceback (most recent call last):\nMemoryError\n", encoding="utf-8"
    )

    healed = resolve(store.read("died"), store)
    assert healed.state == FAILED
    assert "MemoryError" in healed.error


def test_a_dead_job_that_printed_nothing_still_says_it_died(store, tmp_path):
    job = JobStatus(
        id="silent",
        course_dir=str(make_course(tmp_path)),
        state=RUNNING,
        pid=999999,
        started_at=time.time(),
    )
    store.write(job)
    healed = resolve(store.read("silent"), store)
    assert healed.state == FAILED
    assert "printed nothing" in healed.error


def test_resolving_heals_the_record_so_readers_cannot_disagree(store, tmp_path):
    job = JobStatus(
        id="heal",
        course_dir=str(make_course(tmp_path)),
        state=RUNNING,
        pid=999999,
        started_at=time.time(),
    )
    store.write(job)
    resolve(store.read("heal"), store)
    assert store.read("heal").state == FAILED, "the correction was not written back"


def test_a_finished_job_is_never_second_guessed(store, tmp_path):
    job = JobStatus(
        id="done",
        course_dir=str(make_course(tmp_path)),
        state=DONE,
        pid=999999,
        started_at=time.time(),
        finished_at=time.time(),
    )
    store.write(job)
    assert resolve(store.read("done"), store).state == DONE


def test_an_old_record_with_no_pid_falls_back_to_the_timeout(store, tmp_path):
    """Records written before pids were tracked must not all read as dead."""
    job = JobStatus(
        id="old",
        course_dir=str(make_course(tmp_path)),
        state=RUNNING,
        pid=0,
        started_at=time.time(),
    )
    store.write(job)
    assert resolve(store.read("old"), store).state == RUNNING


def test_a_packet_from_before_the_run_started_is_not_this_runs_packet(store, tmp_path):
    """Otherwise every failed re-run inherits the previous run's success."""
    course = make_course(tmp_path)
    wrote_packet(course, time.time() - 3600)

    job = JobStatus(
        id="new", course_dir=str(course), state=RUNNING, pid=999999,
        started_at=time.time(),
    )
    store.write(job)
    assert resolve(store.read("new"), store).state == FAILED


def test_wrote_a_packet_is_false_when_there_is_none(tmp_path):
    assert wrote_a_packet(make_course(tmp_path), time.time()) is False


# ---------------------------------------------------------------------------
# The guard
# ---------------------------------------------------------------------------

def test_a_dead_run_does_not_block_a_new_one(store, tmp_path, monkeypatch):
    """Refusing on a dead job leaves the course unrunnable until a file is deleted."""
    course = make_course(tmp_path)
    store.write(
        JobStatus(
            id="corpse",
            course_dir=str(course),
            state=RUNNING,
            pid=999999,
            started_at=time.time(),
        )
    )
    monkeypatch.setattr("qa.jobs.subprocess.Popen", lambda *a, **k: FakeProcess())
    fresh = submit(course, {}, store)
    assert fresh.id != "corpse"
    assert store.read("corpse").state == FAILED


def test_a_live_run_still_blocks_a_second_one(store, tmp_path, monkeypatch):
    """The guard exists because two runs on one folder invalidate each other."""
    course = make_course(tmp_path)
    store.write(
        JobStatus(
            id="busy",
            course_dir=str(course),
            state=RUNNING,
            pid=os.getpid(),
            started_at=time.time(),
        )
    )
    monkeypatch.setattr("qa.jobs.subprocess.Popen", lambda *a, **k: FakeProcess())
    with pytest.raises(JobError, match="already in progress"):
        submit(course, {}, store)


def test_submit_records_the_pid_and_opens_a_log(store, tmp_path, monkeypatch):
    course = make_course(tmp_path)
    seen: dict = {}

    def fake(*args, **kwargs):
        seen.update(kwargs)
        return FakeProcess(pid=4242)

    monkeypatch.setattr("qa.jobs.subprocess.Popen", fake)
    status = submit(course, {}, store)

    assert status.pid == 4242
    assert store.read(status.id).pid == 4242, "the pid has to reach the record"
    assert seen["stderr"] == subprocess.STDOUT
    assert seen["stdout"] is not subprocess.DEVNULL, "output must be captured"
    assert store.log_path(status.id).exists()


# ---------------------------------------------------------------------------
# What the progress view decides to do
# ---------------------------------------------------------------------------

def test_the_progress_view_keeps_looking_until_a_run_is_finished():
    """The bug, held as a predicate.

    A just-submitted job is PENDING, and the view used to arm its refreshing
    fragment only for RUNNING. So it drew "Queued" once, statically, and never
    looked again while the run went all the way through.
    """
    from qa.web.run_view import should_refresh

    for state in (PENDING, RUNNING):
        assert should_refresh(JobStatus(id="x", course_dir="c", state=state)) is True
    for state in (DONE, FAILED):
        assert should_refresh(JobStatus(id="x", course_dir="c", state=state)) is False


# ---------------------------------------------------------------------------
# End to end, through the same entry point the web layer uses
# ---------------------------------------------------------------------------

def build_runnable_course(root: Path) -> Path:
    """A whole course from nothing: a generated deck and two seconds of tone."""
    import numpy as np
    import soundfile as sf
    from pptx import Presentation

    course = root / "spisccc26" / "course01"
    audio = course / "audio"
    audio.mkdir(parents=True)

    rate = 16000
    time_axis = np.arange(int(2.0 * rate)) / rate
    tone = 0.2 * np.sin(2 * np.pi * 220 * time_axis)
    sf.write(str(audio / "it_gen02_01_enus_01.wav"), tone.astype("float32"), rate)

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Only topic"
    slide.notes_slide.notes_text_frame.text = "The kettle boils a fresh pot."
    deck.save(str(course / "storyboard.pptx"))

    (course / "course.yaml").write_text(
        'course_number: "01"\nproject_type: VENDOR\n'
        "course_code: it_gen02_01_enus\nscript_source: pptx\n"
        "unscripted_topics: []\ntopics: {}\nasr:\n  model: tiny\n",
        encoding="utf-8",
    )
    return course


@pytest.mark.skipif(
    not __import__("qa.setup", fromlist=["check_model"]).check_model("tiny").satisfied,
    reason="no tiny ASR model cached; run qa-setup to fetch one",
)
def test_a_run_started_the_way_the_ui_starts_one_reaches_finished(store, tmp_path):
    """The acceptance case: submit, then watch until the view would say finished.

    Goes through qa.jobs.submit, which is the only way the web layer starts a
    run, and waits on exactly what the progress view reads. It failed before
    this session's fix by never leaving PENDING as far as any reader could see.
    """
    course = build_runnable_course(tmp_path / "library")
    status = submit(
        course,
        {"device": "cpu", "model": "tiny", "output": str(tmp_path / "packets")},
        store,
    )
    assert status.pid, "submit must record the pid it spawned"

    deadline = time.time() + 300
    seen = status
    while time.time() < deadline:
        seen = resolve(store.read(status.id), store)
        if not __import__(
            "qa.web.run_view", fromlist=["should_refresh"]
        ).should_refresh(seen):
            break
        time.sleep(1.0)

    assert seen.state == DONE, f"{seen.state}: {seen.error}\n{store.tail(status.id)}"
    assert seen.topic_total == 1
    packets = list((tmp_path / "packets").glob("*.md"))
    assert len(packets) == 1, "the run wrote no packet"
    assert packets[0].name.startswith("it_gen02_01_enus_")
