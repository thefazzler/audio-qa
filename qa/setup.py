"""Environment check and setup.

    qa-setup            check, then install what is safe to install
    qa-setup --check    the diagnostic table only, changes nothing

Someone who clones this repository should be able to get it running without a
person standing next to them. That is the whole job.

Two kinds of prerequisite, treated differently, per D2:

  Checked and explained, never installed: Python, git, ffmpeg, ffprobe and the
  CUDA runtime. These are system software. The command reports what it found,
  what is required, and the exact command to fix it on this platform, and then
  gets out of the way.

  Installed here, because they are safe and local: the Python packages into the
  project's own virtual environment, and the ASR model. Both live inside the
  project or the user's own cache and neither changes the machine.

This module imports only the standard library at the top level, because it has
to run from a bare checkout before the virtual environment exists. Everything
else is imported inside the function that needs it.
"""

from __future__ import annotations

import argparse
import os
import platform
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path

# Statuses. Kept distinct on purpose: "not installed" and "installed but the
# wrong version" have different fixes and must never collapse into one word.
OK = "OK"
MISSING = "MISSING"
MISMATCH = "VERSION MISMATCH"
NOT_USABLE = "NOT USABLE"
UNKNOWN = "UNKNOWN"

# ctranslate2 does not publish wheels for 3.14, so a newer Python is a trap
# rather than an upgrade. See D9.
MIN_PYTHON = (3, 11)
MAX_PYTHON_EXCLUSIVE = (3, 14)
PREFERRED_PYTHON = "3.12"

# ctranslate2 4.x links against CUDA 12. A machine with CUDA 11 has a real GPU
# and cannot use it, which is a different problem from having no GPU.
CTRANSLATE2_CUDA_MAJOR = 12

MODELS = {
    "large-v3": "about 3 GB, the production model",
    "tiny": "about 75 MB, used by the smoke test",
}

PROJECT_ROOT = Path(__file__).resolve().parent.parent
VENV_DIR = PROJECT_ROOT / ".venv"


@dataclass
class Requirement:
    name: str
    status: str
    found: str = ""
    required: str = ""
    fix: str = ""
    optional: bool = False
    detail: str = ""

    @property
    def satisfied(self) -> bool:
        return self.status == OK

    @property
    def blocking(self) -> bool:
        return not self.satisfied and not self.optional


@dataclass
class SetupReport:
    requirements: list[Requirement] = field(default_factory=list)
    actions: list[str] = field(default_factory=list)
    smoke: str = ""
    smoke_detail: str = ""

    @property
    def blocking(self) -> list[Requirement]:
        return [r for r in self.requirements if r.blocking]

    @property
    def ok(self) -> bool:
        return not self.blocking


# ---------------------------------------------------------------------------
# Platform specific remediation, in one place
# ---------------------------------------------------------------------------

def _system() -> str:
    return platform.system()


def install_hint(tool: str) -> str:
    system = _system()
    hints = {
        "ffmpeg": {
            "Windows": "winget install --id Gyan.FFmpeg -e",
            "Linux": "sudo apt-get update && sudo apt-get install -y ffmpeg",
            "Darwin": "brew install ffmpeg",
        },
        "git": {
            "Windows": "winget install --id Git.Git -e",
            "Linux": "sudo apt-get install -y git",
            "Darwin": "brew install git",
        },
        "python": {
            "Windows": f"winget install --id Python.Python.{PREFERRED_PYTHON} -e",
            "Linux": f"sudo apt-get install -y python{PREFERRED_PYTHON} python{PREFERRED_PYTHON}-venv",
            "Darwin": f"brew install python@{PREFERRED_PYTHON}",
        },
    }
    fallback = {
        "ffmpeg": "install ffmpeg from https://ffmpeg.org/download.html",
        "git": "install git from https://git-scm.com/downloads",
        "python": f"install Python {PREFERRED_PYTHON} from https://python.org",
    }
    return hints.get(tool, {}).get(system) or fallback.get(tool, "")


# ---------------------------------------------------------------------------
# Checks. None of these changes anything.
# ---------------------------------------------------------------------------

def _version_of(command: list[str]) -> str:
    try:
        done = subprocess.run(command, capture_output=True, text=True, timeout=30)
    except (OSError, subprocess.SubprocessError):
        return ""
    text = (done.stdout or done.stderr or "").strip()
    return text.splitlines()[0].strip() if text else ""


def check_python(version: tuple[int, ...] | None = None) -> Requirement:
    """The running interpreter, and whether it is one this project supports."""
    current = version or sys.version_info[:3]
    found = ".".join(str(part) for part in current)
    required = (
        f"{MIN_PYTHON[0]}.{MIN_PYTHON[1]} or newer, "
        f"below {MAX_PYTHON_EXCLUSIVE[0]}.{MAX_PYTHON_EXCLUSIVE[1]}"
    )

    if current[:2] < MIN_PYTHON:
        return Requirement(
            "Python", MISMATCH, found, required,
            f"{install_hint('python')}, then rerun with that interpreter",
        )
    if current[:2] >= MAX_PYTHON_EXCLUSIVE:
        return Requirement(
            "Python", MISMATCH, found, required,
            f"{install_hint('python')}, then rerun with that interpreter",
            detail=(
                f"Python {found} is newer than this project supports. The ASR "
                "runtime does not publish wheels for it, so installing would "
                "fall back to building from source and usually fail. This is a "
                "trap rather than an upgrade; see DECISIONS.md D9."
            ),
        )
    return Requirement("Python", OK, found, required)


def find_supported_python() -> str:
    """A supported interpreter on this machine, for the venv step."""
    if MIN_PYTHON <= sys.version_info[:2] < MAX_PYTHON_EXCLUSIVE:
        return sys.executable

    candidates: list[list[str]] = []
    if _system() == "Windows":
        candidates += [["py", f"-{PREFERRED_PYTHON}"], ["py", "-3.13"], ["py", "-3.11"]]
    candidates += [
        [f"python{PREFERRED_PYTHON}"], ["python3.13"], ["python3.11"], ["python3"]
    ]

    for candidate in candidates:
        executable = shutil.which(candidate[0])
        if not executable:
            continue
        probe = subprocess.run(
            candidate + ["-c", "import sys; print('%d.%d' % sys.version_info[:2])"],
            capture_output=True,
            text=True,
        )
        if probe.returncode != 0:
            continue
        try:
            major, minor = (int(p) for p in probe.stdout.strip().split("."))
        except ValueError:
            continue
        if MIN_PYTHON <= (major, minor) < MAX_PYTHON_EXCLUSIVE:
            return " ".join(candidate)
    return ""


def check_tool(name: str, args: list[str] | None = None) -> Requirement:
    """A system binary that must simply be present and runnable."""
    executable = shutil.which(name)
    if not executable:
        return Requirement(
            name, MISSING, "not on PATH", "any recent version", install_hint(
                "ffmpeg" if name in {"ffmpeg", "ffprobe"} else name
            )
        )
    version = _version_of([executable] + (args or ["--version"]))
    return Requirement(name, OK, version or "present", "any recent version")


def check_git() -> Requirement:
    return check_tool("git")


def check_ffmpeg() -> Requirement:
    return check_tool("ffmpeg")


def check_ffprobe() -> Requirement:
    return check_tool("ffprobe")


# ---------------------------------------------------------------------------
# CUDA. Optional, and the row with the most to say.
# ---------------------------------------------------------------------------

def detect_cuda_major() -> tuple[int | None, str]:
    """Which CUDA major version this machine offers, and where that was seen.

    Looks for a toolkit rather than asking the driver, because the driver's
    "CUDA Version" line reports the newest it could support, not what is
    installed. An absent toolkit is not an error: modern drivers ship the
    runtime components the ASR engine needs.
    """
    for variable in ("CUDA_PATH", "CUDA_HOME"):
        value = os.environ.get(variable)
        if value:
            major = _major_from_path(value)
            if major:
                return major, f"{variable}={value}"

    for entry in os.environ.get("PATH", "").split(os.pathsep):
        if "cuda" not in entry.lower():
            continue
        major = _major_from_path(entry)
        if major:
            return major, f"on PATH: {entry}"
    return None, ""


def _major_from_path(text: str) -> int | None:
    import re

    match = re.search(r"[/\\]v?(\d+)\.(\d+)", text, re.IGNORECASE)
    if match and "cuda" in text.lower():
        return int(match.group(1))
    return None


def cuda_pip_remediation() -> str:
    """Additive, in the virtual environment, changing nothing system wide."""
    from .device import CUDA_PIP_PACKAGES

    return CUDA_PIP_PACKAGES


def check_cuda(probe=None, toolkit=None) -> Requirement:
    """Whether the ASR engine can actually use a GPU on this machine.

    The question is never "is there a card". It is whether ctranslate2 can use
    one, which is a different question and the only one that matters here.

    GPU is optional throughout. The tool runs on CPU without any of this, and
    this row failing does not block setup.
    """
    required = f"CUDA {CTRANSLATE2_CUDA_MAJOR} runtime, usable by ctranslate2"

    if probe is None:
        try:
            from .device import probe_gpu

            probe = probe_gpu()
        except Exception as exc:
            return Requirement(
                "CUDA (GPU)", UNKNOWN, f"probe failed: {exc}", required,
                "GPU is optional; the pipeline runs on CPU.", optional=True,
            )

    major, where = toolkit if toolkit is not None else detect_cuda_major()

    if probe.available:
        found = probe.detail or "usable"
        if major and major != CTRANSLATE2_CUDA_MAJOR:
            found += f"; a CUDA {major} toolkit is also present ({where})"
        return Requirement(
            "CUDA (GPU)", OK, found, required, optional=True,
            detail=(
                "Device may affect decode precision; findings are re-verified. "
                "Measured on Course 11: about 99.4 percent token agreement "
                "between devices. See DECISIONS.md D23."
            ),
        )

    # A card the engine cannot use, with an old toolkit in the way, is a
    # version mismatch and not an absence.
    if major is not None and major < CTRANSLATE2_CUDA_MAJOR:
        return Requirement(
            "CUDA (GPU)",
            MISMATCH,
            f"CUDA {major} found ({where}); ctranslate2 cannot use it",
            required,
            cuda_pip_remediation(),
            optional=True,
            detail=(
                f"ctranslate2 {_ctranslate2_version()} links against CUDA "
                f"{CTRANSLATE2_CUDA_MAJOR}. The pip packages above install the "
                "CUDA 12 runtime and cuDNN into this project's virtual "
                "environment only; they change nothing system wide and are "
                "found ahead of the older libraries because the environment's "
                "own packages are loaded first.\n"
                f"    Do not uninstall the CUDA {major} toolkit. CUDA versions "
                "coexist, and something else on this machine may depend on it. "
                "If you know nothing needs it, removing it is your call and not "
                "this command's business.\n"
                "    GPU is optional: the pipeline runs on CPU without any of "
                "this."
            ),
        )

    reason = getattr(probe, "reason", "") or "no usable CUDA device"
    status = MISSING if "no CUDA capable device" in reason else NOT_USABLE
    return Requirement(
        "CUDA (GPU)", status, reason, required,
        "GPU is optional; the pipeline runs on CPU. To enable it later, "
        f"install a CUDA {CTRANSLATE2_CUDA_MAJOR} capable driver, then "
        + cuda_pip_remediation(),
        optional=True,
    )


def _ctranslate2_version() -> str:
    try:
        import importlib.metadata as metadata

        return metadata.version("ctranslate2")
    except Exception:
        return "4.x"


# ---------------------------------------------------------------------------
# The environment this project installs into
# ---------------------------------------------------------------------------

def venv_python(venv: Path | None = None) -> Path:
    venv = venv or VENV_DIR
    return venv / ("Scripts" if _system() == "Windows" else "bin") / (
        "python.exe" if _system() == "Windows" else "python"
    )


def check_venv(venv: Path | None = None) -> Requirement:
    venv = venv or VENV_DIR
    executable = venv_python(venv)
    if not executable.exists():
        return Requirement(
            "Virtual environment", MISSING, f"no {venv.name} in the project",
            "a venv on a supported Python",
            "qa-setup will create it",
        )
    version = _version_of([str(executable), "--version"])
    return Requirement("Virtual environment", OK, version or str(venv), "a venv")


def check_packages(venv: Path | None = None) -> Requirement:
    """Whether the project and its ASR extra are importable in the venv."""
    executable = venv_python(venv)
    if not executable.exists():
        return Requirement(
            "Python packages", MISSING, "no virtual environment yet",
            'the project plus [asr,dev]', "qa-setup will install them"
        )
    probe = subprocess.run(
        [
            str(executable), "-c",
            "import qa, faster_whisper, soundfile, pptx, docx, yaml, numpy, pytest;"
            "import importlib.metadata as m;"
            "print(m.version('faster-whisper'))",
        ],
        capture_output=True,
        text=True,
    )
    if probe.returncode != 0:
        missing = (probe.stderr or "").strip().splitlines()
        return Requirement(
            "Python packages", MISSING,
            missing[-1] if missing else "not importable",
            "the project plus [asr,dev]",
            'pip install -e ".[asr,dev]"',
        )
    return Requirement(
        "Python packages", OK, f"faster-whisper {probe.stdout.strip()}",
        "the project plus [asr,dev]",
    )


def check_model(name: str = "large-v3") -> Requirement:
    """Whether the ASR model is already in the cache. Downloads nothing."""
    from pathlib import Path as _Path

    home = os.environ.get("HF_HOME") or os.environ.get("HUGGINGFACE_HUB_CACHE")
    roots = [
        _Path(home) if home else _Path.home() / ".cache" / "huggingface",
        _Path.home() / ".cache" / "huggingface" / "hub",
    ]
    needle = f"models--Systran--faster-whisper-{name}"
    for root in roots:
        if not root.exists():
            continue
        for candidate in root.rglob(needle):
            if candidate.is_dir():
                return Requirement(
                    f"ASR model ({name})", OK, f"cached at {candidate.parent}",
                    MODELS.get(name, ""),
                )
    return Requirement(
        f"ASR model ({name})", MISSING, "not downloaded",
        MODELS.get(name, ""),
        "qa-setup offers to download it, with a size warning",
        optional=True,
    )


# ---------------------------------------------------------------------------
# The table
# ---------------------------------------------------------------------------

def check_all(venv: Path | None = None) -> list[Requirement]:
    """Every prerequisite, in the order a person would hit them."""
    return [
        check_python(),
        check_git(),
        check_ffmpeg(),
        check_ffprobe(),
        check_venv(venv),
        check_packages(venv),
        check_model("large-v3"),
        check_cuda(),
    ]


def render_table(requirements: list[Requirement]) -> str:
    """The same output on day one and on the day something breaks."""
    widths = {
        "name": max(len("PREREQUISITE"), *(len(r.name) for r in requirements)),
        "status": max(len("STATUS"), *(len(r.status) for r in requirements)),
    }
    lines = [
        f"  {'PREREQUISITE'.ljust(widths['name'])}  "
        f"{'STATUS'.ljust(widths['status'])}  FOUND",
        "  " + "-" * (widths["name"] + widths["status"] + 60),
    ]
    for item in requirements:
        marker = "" if item.satisfied else (" (optional)" if item.optional else "")
        lines.append(
            f"  {item.name.ljust(widths['name'])}  "
            f"{item.status.ljust(widths['status'])}  {item.found}{marker}"
        )
    return "\n".join(lines)


def render_fixes(requirements: list[Requirement]) -> str:
    problems = [r for r in requirements if not r.satisfied]
    if not problems:
        return ""
    lines = ["", "  What to do:"]
    for item in problems:
        lines.append("")
        label = f"  {item.name}: {item.status}"
        if item.optional:
            label += " (optional, does not block setup)"
        lines.append(label)
        if item.required:
            lines.append(f"    required: {item.required}")
        if item.fix:
            lines.append(f"    fix:      {item.fix}")
        if item.detail:
            for line in item.detail.splitlines():
                lines.append(f"    {line}")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# The things this command does install, because they are safe and local
# ---------------------------------------------------------------------------

def create_venv(venv: Path | None = None, interpreter: str = "") -> tuple[bool, str]:
    """Create the project's virtual environment, if it is not already there."""
    venv = venv or VENV_DIR
    if venv_python(venv).exists():
        return False, f"already present at {venv}"

    interpreter = interpreter or find_supported_python()
    if not interpreter:
        return False, (
            "no supported Python found on this machine. " + install_hint("python")
        )
    done = subprocess.run(
        interpreter.split() + ["-m", "venv", str(venv)],
        capture_output=True,
        text=True,
    )
    if done.returncode != 0:
        return False, f"could not create it: {(done.stderr or '').strip()}"
    return True, f"created at {venv}"


def install_packages(venv: Path | None = None) -> tuple[bool, str]:
    """Install the project and its extras into the project's own environment."""
    executable = venv_python(venv)
    if not executable.exists():
        return False, "no virtual environment to install into"
    if check_packages(venv).satisfied:
        return False, "already installed"

    done = subprocess.run(
        [str(executable), "-m", "pip", "install", "-q", "-e", ".[asr,dev]"],
        cwd=str(PROJECT_ROOT),
        capture_output=True,
        text=True,
    )
    if done.returncode != 0:
        tail = "\n      ".join((done.stderr or "").strip().splitlines()[-6:])
        return False, f"pip failed:\n      {tail}"
    return True, "installed the project and its asr and dev extras"


def download_model(
    name: str = "large-v3", confirm: bool = False, venv: Path | None = None
) -> tuple[bool, str]:
    """Fetch an ASR model, only when a human has said yes.

    Never silent. Today the first transcription downloads three gigabytes with
    no warning, which is a poor thing to discover on a metered connection or
    five minutes before a meeting.
    """
    if check_model(name).satisfied:
        return False, f"{name} is already cached"
    if not confirm:
        return False, (
            f"{name} not downloaded ({MODELS.get(name, '')}). "
            "The first transcription will download it."
        )

    executable = venv_python(venv)
    if not executable.exists():
        return False, "no virtual environment to download into"

    done = subprocess.run(
        [
            str(executable),
            "-c",
            "from faster_whisper import WhisperModel; "
            "WhisperModel(%r, device='cpu', compute_type='int8')" % name,
        ],
        capture_output=True,
        text=True,
    )
    if done.returncode != 0:
        tail = "\n      ".join((done.stderr or "").strip().splitlines()[-4:])
        return False, f"download failed:\n      {tail}"
    return True, f"{name} downloaded and cached"


# ---------------------------------------------------------------------------
# Proof, not a summary
# ---------------------------------------------------------------------------

def build_fixture(target: Path) -> Path:
    """A tiny course made from nothing: a generated clip and a generated deck.

    Generated rather than committed. It contains no course narration, which
    D16 requires, and it keeps binary media out of a repository that has to
    stay free of customer material. The generator is the bundled part.
    """
    import numpy as np
    import soundfile as sf
    from pptx import Presentation

    course = Path(target) / "smoke" / "course01"
    audio = course / "audio"
    audio.mkdir(parents=True, exist_ok=True)

    # Three seconds of a warbling tone. The smoke test proves the chain runs,
    # not that the recogniser is accurate, so the content does not matter.
    rate = 16000
    time = np.arange(int(3.0 * rate)) / rate
    tone = 0.2 * np.sin(2 * np.pi * (180 + 40 * np.sin(2 * np.pi * 0.7 * time)) * time)
    padded = np.concatenate([np.zeros(rate // 2), tone, np.zeros(rate // 2)])
    sf.write(str(audio / "it_smoke_01_enus_01.wav"), padded.astype("float32"), rate)

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Smoke test"
    slide.notes_slide.notes_text_frame.text = (
        "The kettle boils a fresh pot. Volunteers repot the tallest cuttings."
    )
    deck.save(str(course / "it_smoke_01_storyboard.pptx"))

    (course / "course.yaml").write_text(
        'course_number: "01"\n'
        "project_type: VENDOR\n"
        "course_code: it_smoke_01_enus\n"
        "unscripted_topics: []\n"
        "asr:\n"
        "  model: tiny\n",
        encoding="utf-8",
    )
    return course


def smoke_test(workdir: Path | None = None, model: str = "") -> tuple[bool, str]:
    """Run the whole pipeline on the fixture. Installed and works differ.

    Uses the tiny model, so this costs a small download rather than three
    gigabytes, and proves every stage, ffmpeg, soundfile, python-pptx, the ASR
    runtime, alignment and packet writing work together on this machine.
    """
    import contextlib
    import io
    import tempfile

    if not model:
        if check_model("tiny").satisfied:
            model = "tiny"
        elif check_model("large-v3").satisfied:
            model = "large-v3"
        else:
            return False, "skipped: no ASR model is cached, so the chain cannot run"

    created = workdir is None
    root = Path(workdir) if workdir else Path(tempfile.mkdtemp(prefix="qa-smoke-"))
    try:
        course = build_fixture(root)
        from . import cli as cli_module

        cli_module._ASR_OVERRIDES = {"model": model, "cpu_threads": None}
        cli_module._ONLY_TOPICS = None
        cli_module._RUN_DATE = "2026-01-01"

        captured = io.StringIO()
        with contextlib.redirect_stdout(captured):
            cli_module.run(course, None, False)

        packets = list((course / "qa_out").glob("reconciliation_packet_*.md"))
        if not packets:
            return False, "the pipeline ran but produced no packet"
        return True, (
            f"the full pipeline ran on a generated fixture with the {model} "
            f"model and produced {packets[0].name}"
        )
    except Exception as exc:
        return False, f"{type(exc).__name__}: {exc}"
    finally:
        if created:
            shutil.rmtree(root, ignore_errors=True)


# ---------------------------------------------------------------------------
# The command
# ---------------------------------------------------------------------------

def _ask(question: str, assume_yes: bool = False) -> bool:
    if assume_yes:
        print(f"    {question} yes (assumed)")
        return True
    if not sys.stdin or not sys.stdin.isatty():
        print(f"    {question} no (not an interactive terminal)")
        return False
    try:
        answer = input(f"    {question} [y/N] ").strip().lower()
    except (EOFError, KeyboardInterrupt):
        return False
    return answer in {"y", "yes"}


SYSTEM_ITEMS = {"Python", "git", "ffmpeg", "ffprobe"}


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        prog="qa-setup",
        description="Check this machine, then install what is safe to install.",
        epilog=(
            "System software is never installed for you: the check prints the\n"
            "exact command and you run it. Python packages and the ASR model are\n"
            "installed here, because both are local to this project or your own\n"
            "cache. See DECISIONS.md D2 and D22."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--check", action="store_true", help="diagnose only; change nothing"
    )
    parser.add_argument(
        "--yes", action="store_true", help="answer yes to the download prompts"
    )
    parser.add_argument(
        "--no-smoke", action="store_true", help="skip the end to end smoke test"
    )
    parser.add_argument("--model", default="large-v3", help="model to offer to fetch")
    args = parser.parse_args(argv)

    print(f"audio-qa setup: {PROJECT_ROOT}")
    print(f"  {platform.platform()}")
    print()

    report = SetupReport(requirements=check_all())
    print(render_table(report.requirements))
    fixes = render_fixes(report.requirements)
    if fixes:
        print(fixes)

    if args.check:
        print()
        if report.ok:
            print("  Everything required is in place.")
        else:
            print(
                f"  {len(report.blocking)} required item(s) need attention. "
                "Fix them above, then run this again."
            )
        return 0 if report.ok else 1

    system_problems = [r for r in report.blocking if r.name in SYSTEM_ITEMS]
    if system_problems:
        print()
        print(
            "  Stopping. The items above are system software, and this command\n"
            "  does not install system software. Run the commands shown, then\n"
            "  run qa-setup again."
        )
        return 1

    print()
    print("  Setting up:")
    _, message = create_venv()
    print(f"    virtual environment: {message}")
    report.actions.append(message)
    if not venv_python().exists():
        return 1

    _, message = install_packages()
    print(f"    packages: {message}")
    report.actions.append(message)

    if check_model(args.model).satisfied:
        print(f"    model: {args.model} already cached")
    else:
        print()
        print(f"    The {args.model} model is {MODELS.get(args.model, 'large')}.")
        print("    It downloads once and is cached for every course afterwards.")
        if _ask(f"Download {args.model} now?", args.yes):
            _, message = download_model(args.model, confirm=True)
            print(f"    model: {message}")
        else:
            print(
                f"    model: skipped. The first transcription will download "
                f"{args.model} ({MODELS.get(args.model, '')})."
            )

    if not args.no_smoke:
        print()
        if not check_model("tiny").satisfied and not check_model(args.model).satisfied:
            if _ask("Download the tiny model to run the smoke test?", args.yes):
                download_model("tiny", confirm=True)
        print("  Smoke test: running the whole pipeline on a generated fixture...")
        passed, detail = smoke_test()
        report.smoke = "pass" if passed else "fail"
        report.smoke_detail = detail
        print(f"    {'PASS' if passed else 'FAIL'}: {detail}")

    print()
    if report.smoke == "fail":
        print("  Setup finished but the smoke test did not pass. See above.")
        return 1
    print("  Ready. Try:  qa-web    or    qa-run <course_dir>")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
