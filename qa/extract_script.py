"""Module 3: script extraction and topic mapping.

Pulls the narration out of the course's script document, decides which part of
it belongs to which topic, and splits each topic's narration into sentences
with the original text preserved exactly.

The stage is a dispatcher over extractors. A VENDOR course's script is the
speaker notes of a PowerPoint storyboard, which is what this module reads. A
CGT course has no PowerPoint at all and carries its script in a Word document
in the BUS Writing Template, which `qa/extract_docx.py` reads. A topic may also
have a freeform script of its own, or none. Every extractor emits the same
per-topic structure, so the aligner, the checks, the artifacts and the packet
do not know or care which one ran. See `qa/script_source.py` for the
vocabulary and DECISIONS.md D26 for why file type decides none of it.

The topic mapper is the highest-risk component in the pipeline: one slide
assigned to the wrong topic produces a block of false deletions in one topic
and false insertions in the next, which reads exactly like a serious narration
defect. So it does two things beyond mapping. It emits its evidence, naming the
marker phrase that fired on every boundary slide, and it hard-checks its result
against the audio file count from the manifest. A mapper that cannot produce
one topic per delivered audio file stops the run rather than guessing.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from .script_source import (
    DOCX_BUS,
    FREEFORM,
    NONE,
    OUTLINE,
    PPTX,
    VERBATIM,
    TopicScript,
)
from .util import ScriptError, read_json, split_sentences, write_json

# Marker phrases that open a new topic in Skillsoft storyboard notes. Ordered
# most specific first; the first match wins and is recorded as the evidence.
TOPIC_MARKERS: tuple[tuple[str, str], ...] = (
    ("demo_intro", r"^(?:in this (?:demonstration|demo)\b|demonstrate\b)"),
    ("video_intro", r"^in this (?:video|topic)\b"),
    ("summary_intro", r"^in this course,\s*we\b"),
)

_COMPILED_MARKERS = tuple(
    (name, re.compile(pattern, re.IGNORECASE)) for name, pattern in TOPIC_MARKERS
)

# Notes that are template boilerplate rather than narration.
TEMPLATE_HINTS = (
    "general directions",
    "directions for using this template",
    "delete this slide",
    "do not delete",
    "instructions for the",
    "template instructions",
)


@dataclass(frozen=True)
class Slide:
    number: int
    title: str
    notes: str


@dataclass(frozen=True)
class Marker:
    slide: int
    kind: str
    quote: str


def _shape_title(slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return " / ".join(
                line.strip()
                for line in shape.text_frame.text.splitlines()
                if line.strip()
            )
    return ""


def read_slides(storyboard: Path) -> list[Slide]:
    """Every slide in deck order, with its speaker notes verbatim."""
    try:
        deck = Presentation(str(storyboard))
    except Exception as exc:  # python-pptx raises a variety of package errors
        raise ScriptError(f"Could not open storyboard {storyboard.name}:\n  {exc}") from exc

    slides: list[Slide] = []
    for number, slide in enumerate(deck.slides, start=1):
        notes = ""
        if slide.has_notes_slide:
            frame = slide.notes_slide.notes_text_frame
            if frame is not None:
                notes = frame.text or ""
        slides.append(
            Slide(number=number, title=_shape_title(slide), notes=notes)
        )
    if not slides:
        raise ScriptError(f"{storyboard.name} contains no slides.")
    return slides


def exclusion_reason(slide: Slide) -> str | None:
    """Why this slide carries no narration, or None if it does."""
    stripped = slide.notes.strip()
    if not stripped:
        return "empty notes"
    lowered = stripped.lower()
    if any(hint in lowered for hint in TEMPLATE_HINTS):
        return "template instructions"
    return None


def detect_marker(notes: str) -> Marker | None:
    """First narration-continuity marker in a slide's notes, if any."""
    head = notes.strip()
    if not head:
        return None
    # Compare against the opening sentence only: "In this video" three
    # paragraphs down is a back reference, not a topic boundary.
    opening = split_sentences(head)
    probe = opening[0] if opening else head
    for kind, pattern in _COMPILED_MARKERS:
        if pattern.search(probe.strip()):
            return Marker(slide=0, kind=kind, quote=probe.strip()[:120])
    return None


def auto_map(
    slides: list[Slide], topics: list[str]
) -> tuple[dict[str, list[int]], list[Marker], list[dict]]:
    """Map slides to topics by narration continuity.

    Returns (slide_map, markers, excluded). Raises when the number of topics
    the markers imply does not match the number of delivered audio files.
    """
    excluded: list[dict] = []
    narrated: list[Slide] = []
    for slide in slides:
        reason = exclusion_reason(slide)
        if reason:
            excluded.append(
                {"slide": slide.number, "reason": reason, "title": slide.title[:80]}
            )
        else:
            narrated.append(slide)

    if not narrated:
        raise ScriptError("No slide in the storyboard carries speaker notes.")

    markers: list[Marker] = []
    for slide in narrated:
        found = detect_marker(slide.notes)
        if found is not None:
            markers.append(
                Marker(slide=slide.number, kind=found.kind, quote=found.quote)
            )

    # The first narrated block, before any marker, is the course overview and
    # belongs to the first topic.
    boundaries = [m.slide for m in markers]
    lead_in = [s.number for s in narrated if not boundaries or s.number < boundaries[0]]

    starts: list[int] = []
    if lead_in:
        starts.append(lead_in[0])
    starts.extend(boundaries)

    if len(starts) != len(topics):
        raise ScriptError(_mapping_error(starts, markers, topics, lead_in, excluded))

    slide_map: dict[str, list[int]] = {}
    numbers = [s.number for s in narrated]
    for index, topic in enumerate(topics):
        first = starts[index]
        next_start = starts[index + 1] if index + 1 < len(starts) else None
        span = [
            n for n in numbers
            if n >= first and (next_start is None or n < next_start)
        ]
        slide_map[topic] = [span[0], span[-1]]
    return slide_map, markers, excluded


def _mapping_error(
    starts: list[int],
    markers: list[Marker],
    topics: list[str],
    lead_in: list[int],
    excluded: list[dict],
) -> str:
    lines = [
        "PROBABLE MAPPING ERROR: the storyboard implies "
        f"{len(starts)} topics but the course delivered {len(topics)} audio files.",
        "",
        "  Marker evidence:",
    ]
    if lead_in:
        lines.append(
            f"    slide {lead_in[0]:>3}  lead-in       (no marker; opens the first topic)"
        )
    for marker in markers:
        lines.append(f"    slide {marker.slide:>3}  {marker.kind:<13} {marker.quote!r}")
    if not markers:
        lines.append("    (no marker phrase matched any slide)")
    lines += [
        "",
        f"  Audio topics: {', '.join(topics)}",
        f"  Excluded slides: {len(excluded)}",
        "",
        "  Fix one of these:",
        "    - add slide_map to course.yaml to state the mapping explicitly",
        "    - extend TOPIC_MARKERS in qa/extract_script.py if this storyboard",
        "      opens topics with wording the mapper does not know",
        "    - check that every narration file was delivered",
    ]
    return "\n".join(lines)


def _validate_slide_map(
    slide_map: dict[str, list[int]], topics: list[str], slide_count: int
) -> None:
    missing = [t for t in topics if t not in slide_map]
    if missing:
        raise ScriptError(
            "course.yaml slide_map is missing topics: " + ", ".join(missing)
        )
    for topic, (first, last) in slide_map.items():
        if first < 1 or last > slide_count:
            raise ScriptError(
                f"slide_map['{topic}'] = [{first}, {last}] falls outside the "
                f"storyboard, which has {slide_count} slides."
            )


def build_script(
    storyboard: Path,
    topics: list[str],
    unscripted: set[str] | None = None,
    slide_map: dict[str, list[int]] | None = None,
    topic_scripts: dict[str, TopicScript] | None = None,
) -> dict:
    """Extract per-topic narration from the storyboard.

    `topic_scripts` is the general form of `unscripted`, which said only
    outline-or-not. Either may be passed; the explicit states win where both
    speak, and a topic neither mentions is verbatim.
    """
    unscripted = set(unscripted or ())
    states: dict[str, TopicScript] = {t: TopicScript(state=OUTLINE) for t in unscripted}
    states.update(topic_scripts or {})
    slides = read_slides(storyboard)
    by_number = {s.number: s for s in slides}

    if slide_map:
        _validate_slide_map(slide_map, topics, len(slides))
        excluded = [
            {"slide": s.number, "reason": r, "title": s.title[:80]}
            for s in slides
            if (r := exclusion_reason(s))
        ]
        markers = [
            Marker(slide=s.number, kind=m.kind, quote=m.quote)
            for s in slides
            if not exclusion_reason(s) and (m := detect_marker(s.notes))
        ]
        mapping_source = "course.yaml"
        resolved = {t: list(slide_map[t]) for t in topics}
    else:
        resolved, markers, excluded = auto_map(slides, topics)
        mapping_source = "auto"

    excluded_numbers = {e["slide"] for e in excluded}
    topic_entries: list[dict] = []

    for topic in topics:
        first, last = resolved[topic]
        members = [
            by_number[n]
            for n in range(first, last + 1)
            if n in by_number and n not in excluded_numbers
        ]
        sentences: list[str] = []
        sentence_slides: list[int] = []
        for slide in members:
            for sentence in split_sentences(slide.notes):
                sentences.append(sentence)
                sentence_slides.append(slide.number)

        state = states.get(topic, TopicScript(state=VERBATIM))
        span = f"{first}" if first == last else f"{first}-{last}"
        entry: dict = {
            "topic": topic,
            "script": state.state,
            "slides": [first, last],
            "source_ref": f"slides {span}",
            "scripted": state.aligned,
            "slide_numbers": [s.number for s in members],
            "sentences": sentences,
            "sentence_slides": sentence_slides,
            # Slide titles live in the deck's shapes, never in the speaker
            # notes, so a pptx script carries no non-narration text to strip.
            # The key is present so every extractor's output has one shape.
            "non_narration": [],
            "word_count": sum(len(s.split()) for s in sentences),
        }
        if state.state == OUTLINE:
            # Demo outlines are not verbatim narration. Keep the raw text so the
            # packet can show a human what the demo was supposed to cover.
            entry["outline"] = [
                line.strip()
                for slide in members
                for line in re.split(r"[\n\x0b\r]+", slide.notes)
                if line.strip()
            ]
        topic_entries.append(entry)

    return {
        "script_source": PPTX,
        "script_document": storyboard.name,
        "storyboard": storyboard.name,
        "slide_count": len(slides),
        "mapping": {
            "source": mapping_source,
            "markers": [
                {"slide": m.slide, "kind": m.kind, "quote": m.quote} for m in markers
            ],
            "excluded_slides": excluded,
        },
        "warnings": [],
        "topics": topic_entries,
    }


def build_script_for_source(
    source: str,
    document: Path,
    topics: list[str],
    topic_scripts: dict[str, TopicScript],
    course_dir: Path,
    slide_map: dict[str, list[int]] | None = None,
    course_code: str = "",
) -> dict:
    """Dispatch to the extractor for this course's script source.

    The per-topic states that are not the course's own source, `freeform` and
    `none`, are applied afterwards as an overlay, so both course-level
    extractors get them without either knowing about the other.
    """
    if source == PPTX:
        script = build_script(
            storyboard=document,
            topics=topics,
            unscripted=set(),
            slide_map=slide_map,
            topic_scripts=topic_scripts,
        )
    elif source == DOCX_BUS:
        from .extract_docx import build_script_docx_bus

        script = build_script_docx_bus(
            document_path=document,
            topics=topics,
            topic_scripts=topic_scripts,
            course_code=course_code,
        )
    else:
        raise ScriptError(
            f"No extractor for script_source '{source}'. This is a course-level "
            "source, and only pptx and docx_bus are."
        )

    _apply_topic_overlays(script, topic_scripts, course_dir)
    return script


def _apply_topic_overlays(
    script: dict, topic_scripts: dict[str, TopicScript], course_dir: Path
) -> None:
    """Replace what the course document said for topics that overrode it.

    A `freeform` topic's narration is a document of its own, so whatever the
    storyboard had for it is not narration and is discarded. A `none` topic has
    no script at all: its entry keeps its place in the topic map, carries no
    sentences, and the packet shows its transcript instead. Neither is skipped,
    because the audio was still delivered and still has to be checked.
    """
    from .extract_docx import build_freeform_topic

    for entry in script["topics"]:
        state = topic_scripts.get(entry["topic"])
        if state is None:
            continue
        if state.state == FREEFORM:
            replacement = build_freeform_topic(
                entry["topic"], course_dir / state.file
            )
            entry.update(
                {
                    "script": FREEFORM,
                    "scripted": True,
                    # The slides this topic occupies in the deck are no longer
                    # where its script is, so pointing a reader at them would
                    # send them to the wrong document.
                    "slides": None,
                    "source_ref": replacement["source_ref"],
                    "sentences": replacement["sentences"],
                    "non_narration": [],
                    "word_count": replacement["word_count"],
                }
            )
            entry.pop("outline", None)
            _clear_sentence_provenance(entry, drop=True)
        elif state.state == NONE:
            entry.update(
                {
                    "script": NONE,
                    "scripted": False,
                    "sentences": [],
                    "non_narration": [],
                    "word_count": 0,
                }
            )
            entry.pop("outline", None)
            _clear_sentence_provenance(entry, drop=False)


def _clear_sentence_provenance(entry: dict, drop: bool) -> None:
    """Empty the per-sentence back references, without inventing a key.

    Each extractor names this differently, because a storyboard's sentences
    come from slides and a Word script's come from table rows. Only the one
    this entry actually has is touched: adding the other would put a key from
    one extractor's shape into the other's output.
    """
    for key in ("sentence_slides", "sentence_rows"):
        if key not in entry:
            continue
        if drop:
            entry.pop(key)
        else:
            entry[key] = []


def run_extract_script(course_dir: Path, force: bool = False) -> dict:
    """Stage entry point: read manifest.json, write script.json."""
    from .config import load_course_yaml

    cfg = load_course_yaml(course_dir)
    manifest_path = cfg.qa_work / "manifest.json"
    if not manifest_path.exists():
        raise ScriptError(f"{manifest_path} not found. Run the config stage first.")

    manifest = read_json(manifest_path)
    topics = [t["topic"] for t in manifest["topics"]]

    # Record the storyboard the extraction was made from, and say so when it
    # has changed since the last run. The stage runs every time now, so this
    # cannot go stale silently; it is here to make the dependency explicit and
    # to tell the operator that the script they are aligning against is not
    # the one they aligned against yesterday. See DECISIONS.md D17.
    document_sha256 = manifest.get("script_document_sha256") or manifest.get(
        "storyboard_sha256"
    )
    previous_path = cfg.qa_work / "script.json"
    changed = False
    if previous_path.exists() and document_sha256:
        try:
            previous = read_json(previous_path)
            was = previous.get("script_document_sha256") or previous.get(
                "storyboard_sha256"
            )
            changed = bool(was) and was != document_sha256
        except (OSError, ValueError):
            changed = False

    script = build_script_for_source(
        source=cfg.script_source,
        document=cfg.script_document,
        topics=topics,
        topic_scripts=dict(cfg.topic_scripts),
        course_dir=cfg.course_dir,
        slide_map=cfg.slide_map or None,
        course_code=cfg.course_code,
    )
    script["script_document_sha256"] = document_sha256
    script["storyboard_sha256"] = (
        document_sha256 if script.get("storyboard") else None
    )
    if changed:
        script.setdefault("warnings", []).insert(
            0,
            f"{script['script_document']} has changed since the last run. Every "
            "topic is aligned against the new script.",
        )
    write_json(cfg.qa_work / "script.json", script)
    return script
